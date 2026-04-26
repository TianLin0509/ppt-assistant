"""Parse .pptx templates into TemplateMeta with recursive shape traversal."""

from __future__ import annotations

import hashlib
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from src.schema import (
    BBox, FontInfo, ImageSlotInfo, ShapeRole, ShapeType, TemplateMeta,
)


def parse_template(pptx_path: str) -> TemplateMeta:
    path = Path(pptx_path)
    prs = Presentation(str(path))
    slide = prs.slides[0]

    elements: list[ShapeRole] = []
    for z_index, shape in enumerate(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            elements.extend(_parse_group(shape, z_index, None))
        else:
            elements.append(_parse_shape(shape, z_index, None))

    template_id = path.stem.replace(" ", "_")

    return TemplateMeta(
        template_id=template_id,
        file_path=str(path.resolve()),
        file_mtime=path.stat().st_mtime,
        slide_width_emu=prs.slide_width,
        slide_height_emu=prs.slide_height,
        elements=elements,
    )


def _parse_shape(shape, z_index: int, group_path: str | None) -> ShapeRole:
    shape_type = _classify_shape_type(shape)
    bbox = BBox(left=shape.left, top=shape.top, width=shape.width, height=shape.height)

    current_content = None
    first_run_font = None
    paragraph_fonts = None
    paragraph_count = None
    text_hash = None

    if shape.has_text_frame:
        tf = shape.text_frame
        paragraphs = tf.paragraphs
        paragraph_count = len(paragraphs)

        lines = []
        p_fonts = []
        for para in paragraphs:
            lines.append(para.text)
            if para.runs:
                p_fonts.append(_snapshot_font(para.runs[0]))
            else:
                p_fonts.append(FontInfo())

        current_content = "\n".join(lines)
        text_hash = _compute_text_hash(current_content)
        first_run_font = p_fonts[0] if p_fonts else None
        paragraph_fonts = p_fonts if len(p_fonts) > 1 else None

    return ShapeRole(
        shape_id=shape.shape_id,
        shape_name_original=shape.name,
        type=shape_type,
        is_editable=shape_type in (ShapeType.TEXT, ShapeType.IMAGE),
        bbox=bbox,
        text_hash=text_hash,
        current_content=current_content,
        first_run_font=first_run_font,
        paragraph_fonts=paragraph_fonts,
        paragraph_count=paragraph_count,
        is_in_group=group_path is not None,
        group_path=group_path,
        z_order_index=z_index,
    )


def _parse_group(group_shape, z_index: int, parent_path: str | None) -> list[ShapeRole]:
    results = []
    prefix = f"{z_index:02d}" if parent_path is None else f"{parent_path}-{z_index:02d}"
    for i, shape in enumerate(group_shape.shapes):
        child_path = f"{prefix}-{i:02d}"
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            results.extend(_parse_group(shape, i, prefix))
        else:
            results.append(_parse_shape(shape, i, child_path))
    return results


def _classify_shape_type(shape) -> ShapeType:
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.GROUP:
        return ShapeType.GROUP
    if st == MSO_SHAPE_TYPE.TABLE:
        return ShapeType.TABLE
    if st == MSO_SHAPE_TYPE.CHART:
        return ShapeType.CHART
    if st in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
        return ShapeType.IMAGE
    if shape.has_text_frame and shape.text_frame.text.strip():
        return ShapeType.TEXT
    if st == MSO_SHAPE_TYPE.PLACEHOLDER and shape.has_text_frame:
        return ShapeType.TEXT
    return ShapeType.DECORATION


def _snapshot_font(run) -> FontInfo:
    font = run.font
    color_rgb = None
    try:
        if font.color and font.color.type is not None:
            color_rgb = f"#{font.color.rgb}"
    except (AttributeError, TypeError):
        pass

    name_ea = None
    rPr = run._r.find(qn("a:rPr"))
    if rPr is not None:
        ea = rPr.find(qn("a:ea"))
        if ea is not None:
            name_ea = ea.get("typeface")

    size_pt = None
    if font.size is not None:
        size_pt = font.size.pt

    return FontInfo(
        name=font.name,
        name_east_asian=name_ea,
        size_pt=size_pt,
        bold=font.bold,
        italic=font.italic,
        color_rgb=color_rgb,
        underline=font.underline,
    )


def _compute_text_hash(text: str) -> str:
    return hashlib.md5(text.strip().encode("utf-8")).hexdigest()
