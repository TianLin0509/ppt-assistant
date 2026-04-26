"""TD-1 run-level safe text replacement + TD-4 shape multi-location."""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

from src.schema import FontInfo, ShapeRole, ShapeType, TemplateMeta, TextSubtype


def fill_template(
    pptx_path: str,
    meta: TemplateMeta,
    text_choices: dict[str, str],
    output_path: str,
) -> str:
    prs = Presentation(pptx_path)
    slide = prs.slides[0]

    for element in meta.elements:
        if element.type != ShapeType.TEXT or not element.role_key:
            continue
        if element.role_key not in text_choices:
            continue

        shape = _find_shape(slide, element)
        if shape is None:
            continue

        new_text = text_choices[element.role_key]
        subtype = element.text_subtype or TextSubtype.BODY

        if subtype == TextSubtype.BULLET:
            _replace_bullet(
                shape.text_frame, new_text,
                element.paragraph_fonts or ([element.first_run_font] if element.first_run_font else []),
            )
        else:
            _replace_title(shape.text_frame, new_text, element.first_run_font)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def _find_shape(slide, element: ShapeRole):
    suffix = f"[{element.role_key}]"
    for shape in slide.shapes:
        if suffix in shape.name:
            return shape

    for shape in slide.shapes:
        if shape.shape_id == element.shape_id:
            return shape

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if (shape.left == element.bbox.left and shape.top == element.bbox.top
                and shape.width == element.bbox.width):
            return shape

    for shape in slide.shapes:
        if shape.name == element.shape_name_original:
            return shape

    return None


def _replace_title(text_frame, new_text: str, font_info: FontInfo | None) -> None:
    # Clear all paragraphs except the first
    while len(text_frame.paragraphs) > 1:
        p_elem = text_frame.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)

    first_para = text_frame.paragraphs[0]
    # Clear all runs except the first
    _clear_paragraph_runs(first_para)
    if not first_para.runs:
        first_para.add_run()
    first_para.runs[0].text = new_text
    if font_info:
        _apply_font(first_para.runs[0], font_info)


def _replace_bullet(text_frame, new_text: str, para_fonts: list[FontInfo]) -> None:
    lines = new_text.split("\n")
    existing = list(text_frame.paragraphs)
    txBody = text_frame._txBody

    # Fill existing paragraphs
    for i in range(min(len(lines), len(existing))):
        para = existing[i]
        _clear_paragraph_runs(para)
        if not para.runs:
            para.add_run()
        para.runs[0].text = lines[i]
        font = para_fonts[min(i, len(para_fonts) - 1)] if para_fonts else None
        if font:
            _apply_font(para.runs[0], font)

    # Add new paragraphs if more lines than existing
    if len(lines) > len(existing):
        for i in range(len(existing), len(lines)):
            # Deep copy last existing paragraph to preserve formatting (bullet style, indentation)
            template_p = existing[-1]._p
            new_p = deepcopy(template_p)
            # Clear runs in the copy
            for r in new_p.findall(qn("a:r")):
                new_p.remove(r)
            # Add new run with text
            new_r = etree.SubElement(new_p, qn("a:r"))
            new_t = etree.SubElement(new_r, qn("a:t"))
            new_t.text = lines[i]
            txBody.append(new_p)

            # Apply font to new paragraph's run
            font = para_fonts[min(i, len(para_fonts) - 1)] if para_fonts else None
            if font:
                _apply_font(text_frame.paragraphs[-1].runs[0], font)

    # Remove surplus paragraphs
    while len(text_frame.paragraphs) > len(lines):
        p_elem = text_frame.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)


def _clear_paragraph_runs(para) -> None:
    runs = list(para.runs)
    if not runs:
        return
    for run in runs[1:]:
        run._r.getparent().remove(run._r)
    runs[0].text = ""


def _apply_font(run, font_info: FontInfo) -> None:
    font = run.font
    if font_info.name is not None:
        font.name = font_info.name
    if font_info.size_pt is not None:
        font.size = Pt(font_info.size_pt)
    if font_info.bold is not None:
        font.bold = font_info.bold
    if font_info.italic is not None:
        font.italic = font_info.italic
    if font_info.underline is not None:
        font.underline = font_info.underline
    if font_info.color_rgb is not None:
        from pptx.dml.color import RGBColor
        hex_str = font_info.color_rgb.lstrip("#")
        font.color.rgb = RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

    if font_info.name_east_asian:
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font_info.name_east_asian)
