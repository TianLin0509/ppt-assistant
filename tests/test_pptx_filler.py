"""Tests for TD-1 run-level safe text replacement."""

import pytest
from pptx import Presentation
from pptx.util import Pt

from src.core.pptx_parser import parse_template
from src.core.pptx_filler import fill_template
from src.schema import TextSubtype


class TestTitleReplacement:
    def test_title_text_replaced(self, sample_pptx, tmp_path):
        meta = parse_template(str(sample_pptx))
        title = [e for e in meta.elements if e.current_content and "Title" in e.current_content][0]
        title.role_key = "title_main"
        title.text_subtype = TextSubtype.TITLE

        output = tmp_path / "filled.pptx"
        fill_template(str(sample_pptx), meta, {"title_main": "New Title"}, str(output))

        prs = Presentation(str(output))
        slide = prs.slides[0]
        found = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text == "New Title":
                        found = True
                        assert para.runs[0].font.bold is True
                        assert para.runs[0].font.size == Pt(28)
        assert found, "Replaced title not found"

    def test_title_extra_runs_removed(self, sample_pptx, tmp_path):
        meta = parse_template(str(sample_pptx))
        title = [e for e in meta.elements if e.current_content and "Title" in e.current_content][0]
        title.role_key = "title_main"
        title.text_subtype = TextSubtype.TITLE

        output = tmp_path / "filled.pptx"
        fill_template(str(sample_pptx), meta, {"title_main": "Clean"}, str(output))

        prs = Presentation(str(output))
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text == "Clean":
                        assert len(para.runs) == 1


class TestBodyReplacement:
    def test_body_text_replaced(self, sample_pptx, tmp_path):
        meta = parse_template(str(sample_pptx))
        body = [e for e in meta.elements if e.current_content and "Body content" in e.current_content][0]
        body.role_key = "body_left"
        body.text_subtype = TextSubtype.BODY

        output = tmp_path / "filled.pptx"
        fill_template(str(sample_pptx), meta, {"body_left": "Replaced body"}, str(output))

        prs = Presentation(str(output))
        texts = []
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(para.text)
        assert "Replaced body" in texts


class TestBulletReplacement:
    def test_bullet_lines_replaced(self, sample_pptx, tmp_path):
        meta = parse_template(str(sample_pptx))
        bullet = [e for e in meta.elements if e.current_content and "Point one" in e.current_content][0]
        bullet.role_key = "bullet_right"
        bullet.text_subtype = TextSubtype.BULLET

        output = tmp_path / "filled.pptx"
        fill_template(str(sample_pptx), meta, {"bullet_right": "Line A\nLine B\nLine C"}, str(output))

        prs = Presentation(str(output))
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                paras = [p.text for p in shape.text_frame.paragraphs]
                if "Line A" in paras:
                    assert paras == ["Line A", "Line B", "Line C"]
                    assert shape.text_frame.paragraphs[0].runs[0].font.bold is True

    def test_bullet_fewer_lines(self, sample_pptx, tmp_path):
        meta = parse_template(str(sample_pptx))
        bullet = [e for e in meta.elements if e.current_content and "Point one" in e.current_content][0]
        bullet.role_key = "bullet_right"
        bullet.text_subtype = TextSubtype.BULLET

        output = tmp_path / "filled.pptx"
        fill_template(str(sample_pptx), meta, {"bullet_right": "Only one"}, str(output))

        prs = Presentation(str(output))
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                paras = [p.text for p in shape.text_frame.paragraphs if p.text]
                if "Only one" in paras:
                    assert len(paras) == 1


class TestNoMatchSkipped:
    def test_unknown_role_key_skipped(self, sample_pptx, tmp_path):
        meta = parse_template(str(sample_pptx))
        output = tmp_path / "filled.pptx"
        fill_template(str(sample_pptx), meta, {"nonexistent": "value"}, str(output))
        assert output.exists()
