"""Shared test fixtures."""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

SAMPLES_DIR = Path(__file__).parent / "samples"


@pytest.fixture
def sample_pptx(tmp_path) -> Path:
    """Create a minimal .pptx with known shapes for testing."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Shape 0: title text box (top area, large font)
    txbox = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11), Inches(1))
    txbox.name = "Title Box"
    tf = txbox.text_frame
    tf.paragraphs[0].text = ""
    run = tf.paragraphs[0].add_run()
    run.text = "Sample Title"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xC8, 0x31, 0x3A)
    run.font.name = "Arial"

    # Shape 1: body text box (left half, medium font)
    txbox2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(4))
    txbox2.name = "Body Left"
    tf2 = txbox2.text_frame
    tf2.paragraphs[0].text = ""
    run2 = tf2.paragraphs[0].add_run()
    run2.text = "Body content here"
    run2.font.size = Pt(14)
    run2.font.name = "Arial"

    # Shape 2: bullet text box (right half, multi-paragraph)
    txbox3 = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5), Inches(4))
    txbox3.name = "Bullet Right"
    tf3 = txbox3.text_frame
    p1 = tf3.paragraphs[0]
    p1.text = ""
    r1 = p1.add_run()
    r1.text = "Point one"
    r1.font.size = Pt(12)
    r1.font.bold = True
    p2 = tf3.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Point two"
    r2.font.size = Pt(12)
    p3 = tf3.add_paragraph()
    r3 = p3.add_run()
    r3.text = "Point three"
    r3.font.size = Pt(12)

    # Shape 3: image placeholder (add a simple rectangle as stand-in)
    from pptx.enum.shapes import MSO_SHAPE
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(6.5), Inches(2), Inches(0.5))

    out = tmp_path / "test_template.pptx"
    prs.save(str(out))
    return out
